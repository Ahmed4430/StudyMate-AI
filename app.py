import os
import io
import json
import shutil
import hashlib
import fitz
import pytesseract
import streamlit as st
import time
import pymupdf

from pptx import Presentation
from docx import Document
from PIL import Image
from dotenv import load_dotenv
from google import genai
from google.genai import types


# =========================
# PAGE CONFIGURATION
# (must be the first Streamlit command)
# =========================

st.set_page_config(
    page_title="StudyMate AI",
    page_icon="📚",
    layout="wide"
)


# =========================
# TESSERACT OCR
# Cross-platform: only override the tesseract path if a Windows
# install is actually found. On Linux/cloud (e.g. Streamlit Cloud),
# pytesseract will use whatever "tesseract" is on PATH instead.
# =========================

WINDOWS_TESSERACT_PATH = r"C:\Tesse-OCR\tesseract.exe"

if os.path.exists(WINDOWS_TESSERACT_PATH):
    pytesseract.pytesseract.tesseract_cmd = WINDOWS_TESSERACT_PATH
elif shutil.which("tesseract"):
    pytesseract.pytesseract.tesseract_cmd = shutil.which("tesseract")


# =========================
# LOAD API KEY + CACHED CLIENT
# st.cache_resource keeps a single client alive across reruns
# instead of rebuilding it on every interaction.
# =========================

load_dotenv()

MODEL_NAME = "gemini-3.6-flash"


@st.cache_resource(show_spinner=False)
def get_client():
    api_key = os.getenv("GEMINI_API_KEY")
    return genai.Client(api_key=api_key)


client = get_client()


# =========================
# GENERATION CONFIGS
# Gemini 3.x defaults to thinking_level="high", which adds real
# latency to every call. Almost nothing in this app (chat answers,
# summaries, MCQ generation) needs deep reasoning, so we default to
# "low" everywhere. JSON_CONFIG additionally forces clean JSON output
# so quiz/exam generation never needs markdown-fence stripping or
# risks an invalid-JSON retry.
# =========================

FAST_CONFIG = types.GenerateContentConfig(
    thinking_config=types.ThinkingConfig(thinking_level="low"),
)

JSON_CONFIG = types.GenerateContentConfig(
    thinking_config=types.ThinkingConfig(thinking_level="low"),
    response_mime_type="application/json",
)


# =========================
# CUSTOM UI THEME
# Single consolidated style block (the old code injected two
# separate <style> tags). Uses a gradient header, card hover states,
# and CSS variables so it adapts to Streamlit's light/dark theme
# instead of hardcoding a gray background.
# =========================

st.markdown(
    """
    <style>

    .main {
        padding-top: 1rem;
    }

    .study-hero {
        text-align: center;
        padding: 34px 20px 30px 20px;
        border-radius: 18px;
        margin-bottom: 22px;
        background: linear-gradient(135deg, #4F46E5 0%, #7C3AED 50%, #2563EB 100%);
        box-shadow: 0 10px 30px rgba(79, 70, 229, 0.25);
    }

    .study-title {
        font-size: 40px;
        font-weight: 800;
        color: #ffffff;
        margin-bottom: 6px;
        letter-spacing: -0.5px;
    }

    .study-subtitle {
        font-size: 16px;
        color: rgba(255,255,255,0.9);
        font-weight: 500;
    }

    .feature-card {
        padding: 22px 24px;
        border-radius: 16px;
        border: 1px solid rgba(128,128,128,0.18);
        background: rgba(127,127,127,0.05);
        margin-bottom: 18px;
    }

    div[data-testid="stMetric"] {
        background: rgba(127,127,127,0.06);
        border: 1px solid rgba(128,128,128,0.15);
        border-radius: 14px;
        padding: 12px 6px;
    }

    .stButton > button {
        border-radius: 10px;
        font-weight: 600;
        transition: transform 0.12s ease, box-shadow 0.12s ease;
    }

    .stButton > button:hover {
        transform: translateY(-1px);
        box-shadow: 0 4px 12px rgba(79, 70, 229, 0.25);
    }

    .stTabs [data-baseweb="tab-list"] {
        gap: 4px;
    }

    .stTabs [data-baseweb="tab"] {
        font-weight: 600;
        padding: 12px 18px;
        border-radius: 10px 10px 0 0;
    }

    .stTabs [aria-selected="true"] {
        font-weight: 700;
        color: #4F46E5;
    }

    </style>
    """,
    unsafe_allow_html=True
)


# =========================
# MAIN TITLE
# =========================

st.markdown(
    """
    <div class="study-hero">
        <div class="study-title">🤖 StudyMate AI</div>
        <div class="study-subtitle">Your Personal AI Study Agent</div>
    </div>
    """,
    unsafe_allow_html=True
)


# =========================
# WELCOME CARD
# =========================

st.markdown(
    """
    <div class="feature-card">
        <h3>📚 Welcome to StudyMate AI</h3>
        <p>
            Your intelligent study companion for learning,
            revision, quizzes, planning, and exam preparation.
        </p>
    </div>
    """,
    unsafe_allow_html=True
)


# =========================
# FEATURE STATUS
# =========================

col1, col2, col3 = st.columns(3)

with col1:

    st.metric(
        "🤖 AI Tutor",
        "Available"
    )

with col2:

    st.metric(
        "📚 Notes",
        "PDF + OCR"
    )

with col3:

    st.metric(
        "🎯 Exam Mode",
        "Available"
    )
# =========================
# EXAM SESSION STATE
# =========================

if "exam_start_time" not in st.session_state:
    st.session_state.exam_start_time = None

if "exam_time_limit" not in st.session_state:
    st.session_state.exam_time_limit = 30

if "exam_time_expired" not in st.session_state:
    st.session_state.exam_time_expired = False

# =========================
# CREATE TABS
# =========================
chat_tab, planner_tab, notes_tab, quiz_tab, session_tab, progress_tab, exam_tab = st.tabs([
    "💬 AI Tutor",
    "📅 Study Planner",
    "📚 My Notes",
    "🧠 AI Quiz",
    "🎓 Smart Study Session",
    "📊 Progress",
    "🎯 Exam Mode"
])


# =========================
# AI TUTOR WITH CHAT MEMORY
# =========================
with chat_tab:

    st.header("💬 AI Tutor")
    st.caption(
        "Ask questions and get personalized explanations "
        "from StudyMate AI."
    )
    st.write(
        "Ask StudyMate anything about your studies. "
        "You can also ask follow-up questions!"
    )

    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []

    if st.button("🗑️ Clear Chat", key="clear_chat_button"):
        st.session_state.chat_history = []
        st.rerun()

    for message in st.session_state.chat_history:

        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    question = st.chat_input(
        "Ask a study question...",
        key="ai_tutor_input"
    )

    if question:

        st.session_state.chat_history.append({
            "role": "user",
            "content": question
        })

        with st.chat_message("user"):
            st.markdown(question)

        try:

            conversation = """
You are StudyMate AI, a helpful personal AI study tutor.

Answer study-related questions clearly and accurately.
Use the previous conversation to understand follow-up questions.

CONVERSATION:
"""

            for message in st.session_state.chat_history:

                if message["role"] == "user":
                    conversation += f"\nStudent: {message['content']}\n"

                else:
                    conversation += f"\nStudyMate: {message['content']}\n"

            conversation += """
Now respond to the student's latest question.
Give a clear, helpful, and easy-to-understand answer.
"""

            with st.spinner("🤖 StudyMate is thinking..."):

                response = client.models.generate_content(
                    model=MODEL_NAME,
                    contents=conversation,
                    config=FAST_CONFIG
                )

            answer = response.text

            st.session_state.chat_history.append({
                "role": "assistant",
                "content": answer
            })

            with st.chat_message("assistant"):
                st.markdown(answer)

        except Exception as e:
            st.error(f"An error occurred: {e}")


# =========================
# STUDY PLANNER
# =========================
with planner_tab:

    st.header("📅 Create Your Personalized Study Plan")

    st.write(
        "Tell StudyMate about your exams and available time. "
        "AI will create a personalized study schedule for you."
    )

    # -------------------------
    # NORMAL STUDY PLAN
    # -------------------------
    with st.form("study_planner_form"):

        subjects = st.text_area(
            "Subjects",
            placeholder="Example: DBMS, Operating Systems, C Programming"
        )

        days_remaining = st.number_input(
            "Days remaining until exams",
            min_value=1,
            max_value=60,
            value=7
        )

        study_hours = st.number_input(
            "Study hours available per day",
            min_value=1,
            max_value=12,
            value=4
        )

        difficult_subjects = st.text_input(
            "Difficult subjects",
            placeholder="Example: Operating Systems"
        )

        exam_priority = st.text_input(
            "Exam priority",
            placeholder="Example: OS → DBMS → C Programming"
        )

        submitted = st.form_submit_button(
            "✨ Generate My Study Plan"
        )

    if submitted:

        if not subjects.strip():

            st.warning("Please enter at least one subject.")

        else:

            prompt = f"""
You are StudyMate AI, a personal AI study planner.

Create a realistic and personalized study plan.

Student information:
Subjects: {subjects}
Days remaining: {days_remaining}
Study hours available per day: {study_hours}
Difficult subjects: {difficult_subjects}
Exam priority: {exam_priority}

Instructions:
1. Create a plan for every day.
2. Do not exceed the available study hours per day.
3. Give more attention to difficult and high-priority subjects.
4. Include revision and practice where appropriate.
5. Use clear headings: DAY 1, DAY 2, etc.
6. Mention approximate study time for each task.
7. Keep the plan practical and easy to follow.
8. Return only the study plan in a clean format.
"""

            try:

                with st.spinner(
                    "🤖 StudyMate is creating your personalized plan..."
                ):

                    response = client.models.generate_content(
                        model=MODEL_NAME,
                        contents=prompt,
                        config=FAST_CONFIG
                    )

                st.session_state.study_plan = response.text

            except Exception as e:
                st.error(f"An error occurred: {e}")

    # -------------------------
    # ADAPTIVE STUDY PLAN
    # -------------------------
    st.divider()

    st.subheader("🤖 Adaptive Study Plan")

    st.write(
        "Use your quiz results to create a study plan "
        "that gives extra attention to your weak areas."
    )

    if (
        "quiz_history" not in st.session_state
        or len(st.session_state.quiz_history) == 0
    ):

        st.info(
            "Complete at least one quiz first. "
            "Then StudyMate can adapt your study plan."
        )

    else:

        adaptive_days = st.number_input(
            "Days available for your adaptive plan",
            min_value=1,
            max_value=60,
            value=7,
            key="adaptive_days"
        )

        adaptive_hours = st.number_input(
            "Study hours available per day",
            min_value=1,
            max_value=12,
            value=4,
            key="adaptive_hours"
        )

        if st.button(
            "🤖 Create Adaptive Study Plan",
            key="adaptive_plan_button"
        ):

            history_text = ""

            for i, quiz in enumerate(
                st.session_state.quiz_history,
                start=1
            ):

                history_text += (
                    f"Quiz {i}\n"
                    f"Topic: {quiz['topic']}\n"
                    f"Score: {quiz['score']}/{quiz['total']}\n"
                    f"Percentage: {quiz['percentage']}%\n\n"
                )

            adaptive_prompt = f"""
You are StudyMate AI, an intelligent personal study agent.

Create an ADAPTIVE study plan based on the student's quiz performance.

QUIZ PERFORMANCE:
{history_text}

AVAILABLE TIME:
Days available: {adaptive_days}
Study hours per day: {adaptive_hours}

Your task:
1. Analyze which topics have lower scores.
2. Give MORE study time to weaker topics.
3. Give less, but still sufficient, time to stronger topics.
4. Include revision and practice quizzes.
5. Do not exceed {adaptive_hours} hours per day.
6. Create a detailed plan for every day.
7. Use headings: DAY 1, DAY 2, etc.
8. Include the topic and approximate time for every activity.
9. Briefly explain why weak topics received more attention.
10. Base the adaptation ONLY on the quiz performance provided.

Return a clean and practical study plan.
"""

            try:

                with st.spinner(
                    "🤖 StudyMate is analyzing your performance..."
                ):

                    response = client.models.generate_content(
                        model=MODEL_NAME,
                        contents=adaptive_prompt,
                        config=FAST_CONFIG
                    )

                st.session_state.adaptive_study_plan = response.text

            except Exception as e:
                st.error(f"Error creating adaptive study plan: {e}")

    # -------------------------
    # DISPLAY PLANS
    # -------------------------
    if "study_plan" in st.session_state:

        st.divider()
        st.subheader("✨ Your Personalized Study Plan")
        st.markdown(st.session_state.study_plan)

    if "adaptive_study_plan" in st.session_state:

        st.divider()
        st.subheader("🤖 Your Adaptive Study Plan")
        st.markdown(st.session_state.adaptive_study_plan)


# =========================
# MY NOTES / MULTI-FORMAT Q&A
# =========================
with notes_tab:

    st.header("📚 My Notes")

    st.caption(
        "Upload your study materials and use them for "
        "AI-powered questions and summaries."
    )

    st.write(
        "Upload PDF, PowerPoint, Word, text, or image files "
        "and ask questions or generate summaries."
    )

    # =========================
    # FILE UPLOAD
    # =========================

    uploaded_file = st.file_uploader(
        "📂 Upload your study material",
        type=[
            "pdf",
            "pptx",
            "docx",
            "txt",
            "png",
            "jpg",
            "jpeg"
        ],
        key="notes_file_uploader"
    )

    if uploaded_file is not None:

        file_name = uploaded_file.name
        file_extension = file_name.lower().split(".")[-1]
        file_bytes = uploaded_file.getvalue()

        # =========================
        # TEXT EXTRACTION FUNCTIONS
        # =========================

        @st.cache_data(show_spinner=False)
        def extract_pdf_text(file_bytes):

            pdf_document = fitz.open(
                stream=file_bytes,
                filetype="pdf"
            )

            extracted_text = ""
            normal_text_pages = 0
            ocr_pages = 0
            ocr_warnings = []

            for page_number, page in enumerate(pdf_document):

                page_text = page.get_text("text")

                if page_text and page_text.strip():

                    extracted_text += page_text + "\n"
                    normal_text_pages += 1

                else:

                    # OCR FOR SCANNED PDF PAGES
                    try:

                        pix = page.get_pixmap(
                            dpi=200,
                            colorspace=fitz.csRGB,
                            alpha=False
                        )

                        img = Image.frombytes(
                            "RGB",
                            [pix.width, pix.height],
                            pix.samples
                        )

                        ocr_text = pytesseract.image_to_string(
                            img,
                            lang="eng",
                            config="--psm 6"
                        )

                        if ocr_text and ocr_text.strip():

                            extracted_text += ocr_text + "\n"
                            ocr_pages += 1

                        else:

                            ocr_warnings.append(
                                f"⚠️ OCR found no text on page "
                                f"{page_number + 1}"
                            )

                    except Exception as ocr_error:

                        ocr_warnings.append(
                            f"❌ OCR error on page "
                            f"{page_number + 1}: {ocr_error}"
                        )

            return (
                extracted_text,
                normal_text_pages,
                ocr_pages,
                len(pdf_document),
                ocr_warnings
            )

        @st.cache_data(show_spinner=False)
        def extract_pptx_text(file_bytes):

            presentation = Presentation(
                io.BytesIO(file_bytes)
            )

            extracted_text = ""
            slide_count = len(presentation.slides)

            for slide_number, slide in enumerate(
                presentation.slides
            ):

                extracted_text += (
                    f"\n--- Slide {slide_number + 1} ---\n"
                )

                for shape in slide.shapes:

                    if hasattr(shape, "text"):

                        if shape.text.strip():

                            extracted_text += (
                                shape.text + "\n"
                            )

            return extracted_text, slide_count

        @st.cache_data(show_spinner=False)
        def extract_docx_text(file_bytes):

            document = Document(
                io.BytesIO(file_bytes)
            )

            extracted_text = ""

            for paragraph in document.paragraphs:

                if paragraph.text.strip():

                    extracted_text += (
                        paragraph.text + "\n"
                    )

            return extracted_text

        @st.cache_data(show_spinner=False)
        def extract_txt_text(file_bytes):

            try:

                return file_bytes.decode(
                    "utf-8"
                )

            except UnicodeDecodeError:

                return file_bytes.decode(
                    "latin-1"
                )

        @st.cache_data(show_spinner=False)
        def extract_image_text(file_bytes):

            image = Image.open(
                io.BytesIO(file_bytes)
            )

            extracted_text = pytesseract.image_to_string(
                image,
                lang="eng",
                config="--psm 6"
            )

            return extracted_text

        # =========================
        # EXTRACT TEXT BASED ON TYPE
        # =========================

        try:

            extracted_text = ""
            file_info = ""

            # -------------------------
            # PDF
            # -------------------------

            if file_extension == "pdf":

                with st.spinner(
                    "📚 Reading your PDF..."
                ):

                    (
                        extracted_text,
                        normal_text_pages,
                        ocr_pages,
                        total_pages,
                        ocr_warnings
                    ) = extract_pdf_text(
                        file_bytes
                    )

                for warning in ocr_warnings:

                    st.warning(warning)

                file_info = (
                    f"PDF • {total_pages} pages"
                )

                st.success(
                    f"✅ PDF uploaded successfully!"
                )

                st.info(
                    f"Normal text pages: "
                    f"{normal_text_pages} | "
                    f"OCR pages: {ocr_pages} | "
                    f"Extracted characters: "
                    f"{len(extracted_text):,}"
                )

            # -------------------------
            # POWERPOINT
            # -------------------------

            elif file_extension == "pptx":

                with st.spinner(
                    "📊 Reading your PowerPoint..."
                ):

                    (
                        extracted_text,
                        slide_count
                    ) = extract_pptx_text(
                        file_bytes
                    )

                file_info = (
                    f"PowerPoint • "
                    f"{slide_count} slides"
                )

                st.success(
                    "✅ PowerPoint uploaded successfully!"
                )

                st.info(
                    f"Slides: {slide_count} | "
                    f"Extracted characters: "
                    f"{len(extracted_text):,}"
                )

            # -------------------------
            # WORD DOCUMENT
            # -------------------------

            elif file_extension == "docx":

                with st.spinner(
                    "📝 Reading your Word document..."
                ):

                    extracted_text = (
                        extract_docx_text(
                            file_bytes
                        )
                    )

                file_info = "Word Document"

                st.success(
                    "✅ Word document uploaded successfully!"
                )

                st.info(
                    f"Extracted characters: "
                    f"{len(extracted_text):,}"
                )

            # -------------------------
            # TEXT FILE
            # -------------------------

            elif file_extension == "txt":

                with st.spinner(
                    "📄 Reading your text file..."
                ):

                    extracted_text = (
                        extract_txt_text(
                            file_bytes
                        )
                    )

                file_info = "Text File"

                st.success(
                    "✅ Text file uploaded successfully!"
                )

                st.info(
                    f"Extracted characters: "
                    f"{len(extracted_text):,}"
                )

            # -------------------------
            # IMAGE
            # -------------------------

            elif file_extension in [
                "png",
                "jpg",
                "jpeg"
            ]:

                with st.spinner(
                    "🔍 Reading text from your image..."
                ):

                    extracted_text = (
                        extract_image_text(
                            file_bytes
                        )
                    )

                file_info = "Image + OCR"

                st.success(
                    "✅ Image uploaded successfully!"
                )

                st.info(
                    f"OCR extracted characters: "
                    f"{len(extracted_text):,}"
                )

            # =========================
            # NO TEXT FOUND
            # =========================

            if not extracted_text.strip():

                st.error(
                    "❌ StudyMate could not extract text "
                    "from this file."
                )

            else:

                # =========================
                # FILE INFORMATION
                # =========================

                st.caption(
                    f"📎 {file_name} • {file_info}"
                )

                # =========================
                # SUMMARIZE NOTES
                # =========================

                st.subheader(
                    "📝 Summarize Your Notes"
                )

                if st.button(
                    "📝 Summarize Notes",
                    key="summarize_notes_button"
                ):

                    summary_prompt = f"""
You are StudyMate AI, a helpful personal AI study assistant.

Summarize the study material provided below.

Your summary must contain these sections:

## 📖 Overview
Give a short and clear overview of the material.

## 🔑 Main Concepts
List the most important concepts and explain each briefly.

## 📝 Key Points to Remember
Give the important points a student should remember
for studying or exams.

Important rules:
- Use ONLY information provided in the material.
- Do not add information that is not found in the material.
- Keep the summary organized and easy to study.

STUDY MATERIAL:
{extracted_text}
"""

                    try:

                        with st.spinner(
                            "📝 StudyMate is summarizing..."
                        ):

                            response = (
                                client.models.generate_content(
                                    model=MODEL_NAME,
                                    contents=summary_prompt,
                                    config=FAST_CONFIG
                                )
                            )

                        st.session_state.notes_summary = (
                            response.text
                        )

                    except Exception as e:

                        st.error(
                            f"Error summarizing notes: {e}"
                        )

                # =========================
                # DISPLAY SUMMARY
                # =========================

                if "notes_summary" in st.session_state:

                    st.divider()

                    st.subheader(
                        "📝 Your Notes Summary"
                    )

                    st.markdown(
                        st.session_state.notes_summary
                    )

                # =========================
                # ASK QUESTIONS
                # =========================

                st.divider()

                st.subheader(
                    "💬 Ask Questions From Your Notes"
                )

                with st.form(
                    "notes_question_form"
                ):

                    notes_question = st.text_input(
                        "Ask a question about your uploaded material:",
                        placeholder=(
                            "Example: What is a semaphore?"
                        ),
                        key="notes_question_input"
                    )

                    notes_question_submitted = (
                        st.form_submit_button(
                            "🔍 Ask"
                        )
                    )

                if (
                    notes_question_submitted
                    and notes_question.strip()
                ):

                    notes_prompt = f"""
You are StudyMate AI, a helpful personal AI study assistant.

Answer the student's question using ONLY the study material below.

If the answer is not found in the material, clearly say:

"I couldn't find this information in your uploaded notes."

STUDY MATERIAL:
{extracted_text}

STUDENT QUESTION:
{notes_question}

Give a clear and easy-to-understand answer.
"""

                    try:

                        with st.spinner(
                            "📚 Searching your study material..."
                        ):

                            response = (
                                client.models.generate_content(
                                    model=MODEL_NAME,
                                    contents=notes_prompt,
                                    config=FAST_CONFIG
                                )
                            )

                        st.subheader(
                            "🤖 StudyMate's Answer"
                        )

                        st.markdown(
                            response.text
                        )

                    except Exception as e:

                        st.error(
                            f"Error answering question: {e}"
                        )

        except Exception as e:

            st.error(
                f"❌ Error processing {file_name}: {e}"
            )
# =========================
# AI QUIZ
# =========================
with quiz_tab:

    st.header("🎯 AI Quiz Generator")
    st.caption(
        "Test your knowledge with AI-generated multiple-choice questions."
    )

    st.write(
        "Choose a topic, answer the questions, "
        "and let StudyMate evaluate your score."
    )

    quiz_topic = st.text_input(
        "Enter a topic",
        placeholder="Example: DBMS Normalization",
        key="quiz_topic_input"
    )

    num_questions = st.slider(
        "Number of questions",
        min_value=3,
        max_value=10,
        value=5,
        key="quiz_num_input"
    )

    difficulty = st.selectbox(
        "Choose difficulty",
        ["Easy", "Medium", "Hard"],
        key="quiz_difficulty_input"
    )

    # -------------------------
    # GENERATE QUIZ
    # -------------------------
    if st.button(
        "🧠 Generate Quiz",
        key="generate_quiz_button"
    ):

        if not quiz_topic.strip():

            st.warning("Please enter a topic.")

        else:

            quiz_prompt = f"""
You are StudyMate AI, a personal AI study assistant.

Create a multiple-choice quiz.

Topic: {quiz_topic}
Difficulty: {difficulty}
Number of questions: {num_questions}

Return ONLY valid JSON.
Do not use markdown.
Do not add any text before or after the JSON.

Use exactly this structure:

{{
  "questions": [
    {{
      "question": "Question text",
      "options": {{
        "A": "Option A",
        "B": "Option B",
        "C": "Option C",
        "D": "Option D"
      }},
      "correct_answer": "A",
      "explanation": "Brief explanation"
    }}
  ]
}}
"""

            try:

                with st.spinner(
                    "🧠 StudyMate is creating your quiz..."
                ):

                    response = client.models.generate_content(
                        model=MODEL_NAME,
                        contents=quiz_prompt,
                        config=JSON_CONFIG
                    )

                quiz_data = json.loads(response.text)

                st.session_state.quiz_data = quiz_data
                st.session_state.current_quiz_topic = quiz_topic
                st.session_state.quiz_submitted = False
                st.session_state.student_answers = {}

                if "quiz_score" in st.session_state:
                    del st.session_state.quiz_score

                st.success("Quiz generated successfully!")
                st.rerun()

            except json.JSONDecodeError:
                st.error(
                    "The AI generated an invalid quiz format. "
                    "Please generate the quiz again."
                )

            except Exception as e:
                st.error(f"Error generating quiz: {e}")

    # -------------------------
    # DISPLAY QUIZ
    # -------------------------
    if "quiz_data" in st.session_state:

        st.divider()
        st.subheader("📝 Your Quiz")

        questions = st.session_state.quiz_data.get(
            "questions",
            []
        )

        if len(questions) == 0:

            st.error(
                "No questions were found. Please generate the quiz again."
            )

        else:

            for i, question_data in enumerate(questions):

                st.markdown(
                    f"### Question {i + 1}: "
                    f"{question_data.get('question', '')}"
                )

                options = question_data.get("options", {})

                selected_option = st.radio(
                    f"Choose your answer for Question {i + 1}",
                    ["A", "B", "C", "D"],
                    format_func=lambda x, opts=options:
                        f"{x}. {opts.get(x, '')}",
                    index=None,
                    key=f"quiz_answer_{i}"
                )

                st.session_state.student_answers[i] = selected_option

            # -------------------------
            # SUBMIT QUIZ
            # -------------------------
            if not st.session_state.quiz_submitted:

                if st.button(
                    "📤 Submit Quiz",
                    key="submit_quiz_button"
                ):

                    unanswered = [
                        i + 1
                        for i in range(len(questions))
                        if st.session_state.student_answers.get(i) is None
                    ]

                    if unanswered:

                        st.warning(
                            "Please answer all questions first. "
                            f"Unanswered: "
                            f"{', '.join(map(str, unanswered))}"
                        )

                    else:

                        score = 0

                        for i, question_data in enumerate(questions):

                            if (
                                st.session_state.student_answers.get(i)
                                == question_data.get("correct_answer")
                            ):
                                score += 1

                        total = len(questions)
                        percentage = round((score / total) * 100)

                        st.session_state.quiz_score = score
                        st.session_state.quiz_percentage = percentage
                        st.session_state.quiz_submitted = True

                        if "quiz_history" not in st.session_state:
                            st.session_state.quiz_history = []

                        st.session_state.quiz_history.append({
                            "topic": st.session_state.current_quiz_topic,
                            "score": score,
                            "total": total,
                            "percentage": percentage
                        })

                        st.rerun()

            # -------------------------
            # QUIZ RESULTS
            # -------------------------
            if st.session_state.quiz_submitted:

                score = st.session_state.quiz_score
                total = len(questions)
                percentage = st.session_state.quiz_percentage

                st.divider()
                st.header("📊 Quiz Results")

                st.success(
                    f"Your Score: {score}/{total} ({percentage}%)"
                )

                for i, question_data in enumerate(questions):

                    student_answer = (
                        st.session_state.student_answers.get(i)
                    )

                    correct_answer = question_data.get(
                        "correct_answer"
                    )

                    st.markdown(f"### Question {i + 1}")

                    if student_answer == correct_answer:

                        st.success(
                            f"✅ Correct! Your answer: {student_answer}"
                        )

                    else:

                        st.error(
                            f"❌ Your answer: {student_answer}"
                        )

                        st.info(
                            f"Correct answer: {correct_answer}"
                        )

                    with st.expander(
                        f"📖 Show Explanation for Question {i + 1}"
                    ):

                        st.write(
                            question_data.get(
                                "explanation",
                                "No explanation available."
                            )
                        )


# =========================
# STUDY PROGRESS + AI ANALYSIS
# =========================
with progress_tab:

    st.header("📊 My Study Progress")
    st.caption(
        "Track your quiz, study session, and exam performance over time."
   )

    if (
        "quiz_history" not in st.session_state
        or len(st.session_state.quiz_history) == 0
    ):

        st.info(
            "No quiz results yet. Complete an AI Quiz "
            "to start tracking your progress!"
        )

    else:

        history = st.session_state.quiz_history

        total_quizzes = len(history)

        average_score = sum(
            quiz["percentage"] for quiz in history
        ) / total_quizzes

        best_score = max(
            quiz["percentage"] for quiz in history
        )

        col1, col2, col3 = st.columns(3)

        with col1:
            st.metric(
                "🎯 Quizzes Completed",
                total_quizzes
            )

        with col2:
            st.metric(
                "📈 Average Score",
                f"{average_score:.0f}%"
            )

        with col3:
            st.metric(
                "🏆 Best Score",
                f"{best_score}%"
            )

        st.divider()

        st.subheader("📚 Quiz History")

        for i, quiz in enumerate(
            reversed(history),
            start=1
        ):

            st.write(
                f"**{i}. {quiz['topic']}** — "
                f"Score: {quiz['score']}/{quiz['total']} "
                f"({quiz['percentage']}%)"
            )

        st.divider()

        # -------------------------
        # AI PERFORMANCE ANALYSIS
        # -------------------------
        st.subheader("🤖 AI Performance Analysis")

        if st.button(
            "🧠 Analyze My Performance",
            key="analyze_performance_button"
        ):

            history_text = ""

            for i, quiz in enumerate(history, start=1):

                history_text += (
                    f"Quiz {i}\n"
                    f"Topic: {quiz['topic']}\n"
                    f"Score: {quiz['score']}/{quiz['total']}\n"
                    f"Percentage: {quiz['percentage']}%\n\n"
                )

            analysis_prompt = f"""
You are StudyMate AI, a personal AI study agent.

Analyze the student's quiz performance below.

QUIZ HISTORY:
{history_text}

Give a personalized analysis using exactly these sections:

## 📊 Overall Performance

## 💪 Strong Areas

## ⚠️ Areas to Improve

## 📚 Smart Recommendations

Give exactly 3 practical recommendations.

## 🎯 Recommended Next Step

Suggest exactly what the student should study or practice next.

Important:
Base your analysis ONLY on the quiz history.
If there is not enough data, say more quizzes are needed.
Keep the response clear and encouraging.
"""

            try:

                with st.spinner(
                    "🤖 StudyMate is analyzing your performance..."
                ):

                    response = client.models.generate_content(
                        model=MODEL_NAME,
                        contents=analysis_prompt,
                        config=FAST_CONFIG
                    )

                st.session_state.performance_analysis = response.text

            except Exception as e:
                st.error(f"Error analyzing performance: {e}")

        if "performance_analysis" in st.session_state:

            st.success(
                "Your personalized performance analysis is ready!"
            )

            st.markdown(
                st.session_state.performance_analysis
            )

        st.divider()

        if st.button(
            "🗑️ Clear Progress History",
            key="clear_progress_button"
        ):

            st.session_state.quiz_history = []

            if "performance_analysis" in st.session_state:
                del st.session_state.performance_analysis

            st.success("Progress history cleared!")
            st.rerun()


# =========================
# SMART STUDY SESSION
# =========================
with session_tab:

    st.header("🧠 Smart Study Session")
    st.caption(
        "Follow a focused AI-powered study session based on your topic and difficulty."
)

    st.write(
        "Learn a topic, take a quick quiz, "
        "and get a personalized recommendation."
    )

    session_topic = st.text_input(
        "What do you want to study?",
        placeholder="Example: DBMS Normalization",
        key="session_topic_input"
    )

    session_difficulty = st.selectbox(
        "Choose your level",
        ["Beginner", "Intermediate", "Advanced"],
        key="session_difficulty"
    )

    # -------------------------
    # START SESSION
    # -------------------------
    if st.button(
        "🚀 Start Study Session",
        key="start_study_session"
    ):

        if not session_topic.strip():

            st.warning("Please enter a topic.")

        else:

            learning_prompt = f"""
You are StudyMate AI, a personal AI study tutor.

Teach the following topic to a student.

Topic: {session_topic}
Student level: {session_difficulty}

Your explanation should include:

1. What the topic means
2. Important concepts
3. A simple example
4. Key points to remember

Make the explanation clear and suitable for the student's level.
"""

            try:

                with st.spinner(
                    "📖 StudyMate is preparing your lesson..."
                ):

                    response = client.models.generate_content(
                        model=MODEL_NAME,
                        contents=learning_prompt,
                        config=FAST_CONFIG
                    )

                st.session_state.current_session_topic = session_topic
                st.session_state.current_session_level = (
                    session_difficulty
                )
                st.session_state.session_lesson = response.text
                st.session_state.session_quiz_data = None
                st.session_state.session_answers = {}
                st.session_state.session_submitted = False

                st.success("Your study session is ready!")
                st.rerun()

            except Exception as e:
                st.error(f"Error starting study session: {e}")

    # -------------------------
    # DISPLAY LESSON
    # -------------------------
    if "session_lesson" in st.session_state:

        st.divider()
        st.subheader("📖 Learn")

        st.markdown(st.session_state.session_lesson)

        # -------------------------
        # GENERATE QUICK QUIZ
        # -------------------------
        if st.button(
            "🎯 Take a Quick Quiz",
            key="take_session_quiz"
        ):

            quiz_prompt = f"""
You are StudyMate AI.

Create a quick multiple-choice quiz.

Topic: {st.session_state.current_session_topic}
Student level: {st.session_state.current_session_level}

Create exactly 3 questions.

Return ONLY valid JSON.
Do not use markdown.
Do not add text before or after the JSON.

Use exactly this structure:

{{
  "questions": [
    {{
      "question": "Question text",
      "options": {{
        "A": "Option A",
        "B": "Option B",
        "C": "Option C",
        "D": "Option D"
      }},
      "correct_answer": "A",
      "explanation": "Brief explanation"
    }}
  ]
}}
"""

            try:

                with st.spinner(
                    "🧠 Creating your quick quiz..."
                ):

                    response = client.models.generate_content(
                        model=MODEL_NAME,
                        contents=quiz_prompt,
                        config=JSON_CONFIG
                    )

                session_quiz = json.loads(response.text)

                st.session_state.session_quiz_data = session_quiz
                st.session_state.session_answers = {}
                st.session_state.session_submitted = False

                st.rerun()

            except json.JSONDecodeError:
                st.error(
                    "The quiz format was invalid. Please try again."
                )

            except Exception as e:
                st.error(f"Error creating quiz: {e}")

    # -------------------------
    # DISPLAY SESSION QUIZ
    # -------------------------
    if st.session_state.get("session_quiz_data"):

        st.divider()
        st.subheader("🎯 Quick Quiz")

        session_questions = (
            st.session_state.session_quiz_data.get(
                "questions",
                []
            )
        )

        for i, question_data in enumerate(session_questions):

            st.markdown(
                f"### Question {i + 1}: "
                f"{question_data.get('question', '')}"
            )

            options = question_data.get("options", {})

            selected_answer = st.radio(
                f"Choose your answer for Question {i + 1}",
                ["A", "B", "C", "D"],
                format_func=lambda x, opts=options:
                    f"{x}. {opts.get(x, '')}",
                index=None,
                key=f"session_answer_{i}"
            )

            st.session_state.session_answers[i] = selected_answer

        if not st.session_state.get("session_submitted", False):

            if st.button(
                "📤 Submit Quick Quiz",
                key="submit_session_quiz"
            ):

                unanswered = [
                    i + 1
                    for i in range(len(session_questions))
                    if st.session_state.session_answers.get(i) is None
                ]

                if unanswered:

                    st.warning(
                        "Please answer all questions first. "
                        f"Unanswered: "
                        f"{', '.join(map(str, unanswered))}"
                    )

                else:

                    score = 0

                    for i, question_data in enumerate(
                        session_questions
                    ):

                        if (
                            st.session_state.session_answers.get(i)
                            == question_data.get("correct_answer")
                        ):
                            score += 1

                    total = len(session_questions)
                    percentage = round((score / total) * 100)

                    st.session_state.session_score = score
                    st.session_state.session_percentage = percentage
                    st.session_state.session_submitted = True

                    st.rerun()

    # -------------------------
    # SESSION RESULTS
    # -------------------------
    if (
        st.session_state.get("session_quiz_data")
        and st.session_state.get("session_submitted")
    ):

        st.divider()
        st.subheader("📊 Session Results")

        score = st.session_state.session_score
        percentage = st.session_state.session_percentage

        st.success(
            f"Your Score: {score}/3 ({percentage}%)"
        )

        if percentage >= 80:

            recommendation = (
                "Excellent! You understand this topic well. "
                "Try a harder quiz or move to a related topic."
            )

        elif percentage >= 50:

            recommendation = (
                "Good progress! Review the important concepts and "
                "take another quiz for more practice."
            )

        else:

            recommendation = (
                "You may need more practice. Review the lesson carefully, "
                "focus on the concepts you missed, and try again."
            )

        st.info(
            f"🤖 StudyMate Recommendation: {recommendation}"
        )

        st.subheader("📖 Review Your Answers")

        session_questions = (
            st.session_state.session_quiz_data.get(
                "questions",
                []
            )
        )

        for i, question_data in enumerate(session_questions):

            student_answer = (
                st.session_state.session_answers.get(i)
            )

            correct_answer = question_data.get(
                "correct_answer"
            )

            st.markdown(f"### Question {i + 1}")

            if student_answer == correct_answer:

                st.success(
                    f"✅ Correct! Answer: {student_answer}"
                )

            else:

                st.error(
                    f"❌ Your answer: {student_answer}"
                )

                st.info(
                    f"Correct answer: {correct_answer}"
                )

            with st.expander(
                f"📖 Explanation for Question {i + 1}"
            ):

                st.write(
                    question_data.get(
                        "explanation",
                        "No explanation available."
                    )
                )

# =========================
# EXAM MODE
# =========================
with exam_tab:

    st.header("🎯 Exam Mode")
    st.caption(
       "Practice with AI-generated exams and review your performance."
    )

    st.write(
        "Practice for your exams with personalized "
        "multiple-choice questions."
    )

    # =========================
    # EXAM SETTINGS
    # =========================

    exam_subject = st.text_input(
        "📚 Subject / Topic",
        placeholder="Example: Operating Systems",
        key="exam_subject_input"
    )

    exam_difficulty = st.selectbox(
        "🎚️ Difficulty",
        ["Easy", "Medium", "Hard"],
        key="exam_difficulty_input"
    )

    exam_question_count = st.selectbox(
        "🔢 Number of Questions",
        [5, 10, 15],
        key="exam_question_count_input"
    )

    # =========================
    # GENERATE EXAM
    # =========================

    if st.button(
        "🚀 Generate Exam",
        key="generate_exam_button"
    ):

        if not exam_subject.strip():

            st.warning(
                "Please enter a subject or topic first."
            )

        else:

            exam_prompt = f"""
You are StudyMate AI, an AI exam preparation assistant.

Create a practice exam.

Subject:
{exam_subject}

Difficulty:
{exam_difficulty}

Number of questions:
{exam_question_count}

Create exactly {exam_question_count}
multiple-choice questions.

Return ONLY valid JSON.

Use exactly this structure:

{{
  "questions": [
    {{
      "question": "Question text",
      "options": {{
        "A": "Option A",
        "B": "Option B",
        "C": "Option C",
        "D": "Option D"
      }},
      "correct_answer": "A",
      "explanation": "Short explanation"
    }}
  ]
}}
"""

            try:

                with st.spinner(
                    "🤖 StudyMate is creating your exam..."
                ):

                    response = client.models.generate_content(
                        model=MODEL_NAME,
                        contents=exam_prompt,
                        config=JSON_CONFIG
                    )

                exam_text = response.text.strip()

                if exam_text.startswith("```"):

                    exam_text = (
                        exam_text
                        .replace("```json", "")
                        .replace("```", "")
                        .strip()
                    )

                exam_data = json.loads(exam_text)

                # Save exam
                st.session_state.exam_data = exam_data
                st.session_state.exam_answers = {}
                st.session_state.exam_submitted = False
                st.session_state.exam_score = 0
                st.session_state.exam_time_expired = False

                # Start timer: ~90 seconds per question, so a 10-question
                # exam gets 15 minutes. (Previously hardcoded to a flat
                # 30 seconds regardless of exam length or the "10-minute"
                # comment here — that was a bug.)
                st.session_state.exam_start_time = time.time()
                st.session_state.exam_time_limit = exam_question_count * 90

                st.rerun()

            except json.JSONDecodeError:

                st.error(
                    "The AI returned an invalid exam format. "
                    "Please try again."
                )

            except Exception as e:

                st.error(
                    f"Error creating exam: {e}"
                )

    # =========================
    # DISPLAY EXAM
    # =========================

    if st.session_state.get("exam_data"):

        questions = st.session_state.exam_data.get(
            "questions",
            []
        )

        def score_exam(questions):
            """Shared scoring logic used by both the manual Submit
            button and the auto-timeout path below, so an exam that
            times out is graded exactly the same way as one that's
            submitted manually (unanswered questions just count as
            wrong)."""

            score = 0

            for i, question_data in enumerate(questions):

                if (
                    st.session_state.exam_answers.get(i)
                    == question_data.get("correct_answer")
                ):
                    score += 1

            st.session_state.exam_score = score
            st.session_state.exam_submitted = True

        # =========================
        # LIVE TIMER (auto-submits on expiry)
        #
        # st.fragment(run_every=1) reruns just this small piece of the
        # page once a second WITHOUT reloading the browser or losing
        # session_state (unlike the old approach, which only animated a
        # countdown in JavaScript and never actually told Streamlit
        # time was up unless the student clicked something).
        #
        # Requires Streamlit >= 1.37. If your deployed version is older,
        # upgrade streamlit in requirements.txt.
        # =========================

        @st.fragment(run_every=1)
        def exam_timer():

            if st.session_state.get("exam_submitted"):
                return

            if st.session_state.get("exam_start_time") is not None:

                elapsed_time = (
                    time.time() - st.session_state.exam_start_time
                )

                remaining_time = (
                    st.session_state.exam_time_limit - elapsed_time
                )

            else:

                remaining_time = st.session_state.get(
                    "exam_time_limit", 10 * 60
                )

            if remaining_time <= 0:

                # Time's up — grade whatever was answered so far and
                # switch the whole page over to the results view.
                score_exam(questions)
                st.session_state.exam_time_expired = True
                st.rerun()

            else:

                remaining_seconds = max(0, int(remaining_time))
                minutes, seconds = divmod(remaining_seconds, 60)

                urgent = remaining_seconds <= 30
                color = "#DC2626" if urgent else "inherit"

                st.markdown(
                    f"""
                    <div style="
                        text-align:center;
                        font-size:28px;
                        font-weight:bold;
                        padding:15px;
                        color:{color};
                    ">
                        ⏱️ Time Remaining:
                        {minutes:02d}:{seconds:02d}
                    </div>
                    """,
                    unsafe_allow_html=True
                )

        if not st.session_state.get("exam_submitted"):
            exam_timer()

        elif st.session_state.get("exam_time_expired"):

            st.warning(
                "⏰ Time is up! Your exam has been submitted."
            )

        st.divider()

        st.subheader(
            "📝 Your Practice Exam"
        )

        # =========================
        # QUESTIONS
        # =========================

        if not st.session_state.get(
            "exam_submitted",
            False
        ):

            for i, question_data in enumerate(
                questions
            ):

                st.markdown(
                    f"### Question {i + 1}"
                )

                st.write(
                    question_data.get(
                        "question",
                        ""
                    )
                )

                options = question_data.get(
                    "options",
                    {}
                )

                selected_answer = st.radio(
                    "Choose your answer:",
                    ["A", "B", "C", "D"],
                    format_func=lambda x, opts=options:
                        f"{x}) {opts.get(x, '')}",
                    index=None,
                    key=f"exam_answer_{i}"
                )

                st.session_state.exam_answers[i] = (
                    selected_answer
                )

            # =========================
            # SUBMIT EXAM
            # =========================

            if st.button(
                "📤 Submit Exam",
                key="submit_exam_button"
            ):

                unanswered = [
                    i + 1
                    for i in range(len(questions))
                    if st.session_state.exam_answers.get(i)
                    is None
                ]

                if unanswered:

                    st.warning(
                        "Please answer all questions first. "
                        f"Unanswered: "
                        f"{', '.join(map(str, unanswered))}"
                    )

                else:

                    score_exam(questions)

                    st.rerun()

        # =========================
        # EXAM RESULTS
        # =========================

        if st.session_state.get(
            "exam_submitted",
            False
        ):

            st.divider()

            st.subheader(
                "📊 Exam Results"
            )

            score = st.session_state.exam_score
            total = len(questions)

            percentage = (
                round((score / total) * 100)
                if total
                else 0
            )

            st.success(
                f"🎯 Your Score: "
                f"{score}/{total} "
                f"({percentage}%)"
            )

            if percentage >= 80:

                st.info(
                    "🌟 Excellent! You have a strong "
                    "understanding of this topic."
                )

            elif percentage >= 50:

                st.info(
                    "👍 Good progress! Review the "
                    "questions you missed and practice again."
                )

            else:

                st.warning(
                    "📚 Keep practicing! Review the topic "
                    "and try the exam again."
                )

            # =========================
            # REVIEW ANSWERS
            # =========================

            st.subheader(
                "📖 Review Your Answers"
            )

            for i, question_data in enumerate(
                questions
            ):

                student_answer = (
                    st.session_state.exam_answers.get(i)
                )

                correct_answer = (
                    question_data.get(
                        "correct_answer"
                    )
                )

                st.markdown(
                    f"### Question {i + 1}"
                )

                if student_answer == correct_answer:

                    st.success(
                        f"✅ Correct! "
                        f"Your answer: {student_answer}"
                    )

                else:

                    st.error(
                        f"❌ Your answer: "
                        f"{student_answer}"
                    )

                    st.info(
                        f"Correct answer: "
                        f"{correct_answer}"
                    )

                with st.expander(
                    f"📖 Explanation {i + 1}"
                ):

                    st.write(
                        question_data.get(
                            "explanation",
                            "No explanation available."
                        )
                    )
